#!/usr/bin/env python3
"""Generate MaintainX Pro hackathon presentation."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "MaintainX-Pro-Hackathon.pptx"

# Brand palette
BRAND = RGBColor(0x4D, 0x1B, 0xFF)
BRAND_LIGHT = RGBColor(0x6B, 0x4D, 0xFF)
ACCENT = RGBColor(0x15, 0x9B, 0xB6)
DARK = RGBColor(0x0F, 0x0F, 0x12)
INK = RGBColor(0x1D, 0x1D, 0x1F)
MUTED = RGBColor(0x6E, 0x6E, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.transparency = alpha
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=INK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=INK, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = Pt(spacing)
        p.bullet = True
    return box


def slide_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), BRAND)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.7),
                title, size=32, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.5),
                    subtitle, size=14, color=MUTED)


def slide_title_dark(prs, title, subtitle, tagline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BRAND)
    add_rect(slide, Inches(10.5), Inches(-1), Inches(4), Inches(4), BRAND, alpha=0.85)
    add_rect(slide, Inches(-1), Inches(5), Inches(4), Inches(4), ACCENT, alpha=0.9)
    add_textbox(slide, Inches(0.9), Inches(2.2), Inches(10), Inches(1.2),
                title, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.9), Inches(3.4), Inches(10), Inches(0.8),
                subtitle, size=22, color=RGBColor(0xAE, 0xAE, 0xB2))
    if tagline:
        add_textbox(slide, Inches(0.9), Inches(5.8), Inches(10), Inches(0.5),
                    tagline, size=13, color=RGBColor(0x86, 0x86, 0x8B))
    return slide


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BRAND)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(1),
                "Merci !", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(3.6), Inches(11), Inches(0.6),
                "Questions ?", size=28, color=RGBColor(0xE0, 0xDC, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.8),
                "maintanxpro.netlify.app  ·  admin@maintainx.com / demo1234",
                size=16, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 1. TITLE ──
    slide_title_dark(
        prs,
        "MaintainX Pro",
        "Smart Industrial Maintenance Platform",
        "Smart Industrial Asset Intelligence & Spare Parts Innovation Hackathon · 2026",
    )
    s = prs.slides[-1]
    add_textbox(s, Inches(0.9), Inches(4.2), Inches(10), Inches(0.5),
                "De la panne signalée à la pièce commandée — en un seul flux.",
                size=18, color=WHITE)

    # ── 2. PROBLÈME ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Le constat terrain", "Pourquoi la maintenance industrielle mérite mieux")
    add_bullets(slide, Inches(0.9), Inches(1.6), Inches(5.8), Inches(4.5), [
        "Pannes signalées par téléphone, WhatsApp ou papier → perte d'information",
        "Aucune preuve visuelle systématique → diagnostics plus lents et litiges",
        "Workflow opaque : qui fait quoi, où en est la pièce détachée ?",
        "Managers sans vision temps réel du parc et des risques",
        "Clients industriels exclus du suivi de leurs propres équipements",
    ], size=17)
    add_rect(slide, Inches(7.2), Inches(1.6), Inches(5.2), Inches(4.8), RGBColor(0xF5, 0xF5, 0xF7))
    add_textbox(slide, Inches(7.5), Inches(1.9), Inches(4.6), Inches(0.5),
                "Impact chiffré (industrie)", size=16, bold=True, color=BRAND)
    stats = [
        ("−30 à 45 min", "perdues par panne\n(saisie + recherche d'info)"),
        ("×3", "délai si pièce\nmal référencée"),
        ("60 %", "des arrêts évitables\navec maintenance prédictive"),
    ]
    y = 2.5
    for val, lbl in stats:
        add_textbox(slide, Inches(7.5), Inches(y), Inches(2), Inches(0.6),
                    val, size=28, bold=True, color=INK)
        add_textbox(slide, Inches(9.5), Inches(y), Inches(2.8), Inches(0.8),
                    lbl, size=13, color=MUTED)
        y += 1.35

    # ── 3. SOLUTION ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Notre réponse", "Une plateforme PWA tout-en-un")
    pillars = [
        ("📸", "Capture visuelle guidée", "Scan multi-angles, preuves obligatoires, pré-diagnostic assisté"),
        ("📋", "Workflow Kanban", "7 étapes de la panne à la clôture, drag & drop intuitif"),
        ("🔧", "Pièces détachées", "Commande, réception, installation — liées à chaque panne"),
        ("📊", "Intelligence opérationnelle", "Health score, dashboard KPI, jumeau 2D/3D des chantiers"),
    ]
    x = 0.7
    for icon, title, desc in pillars:
        add_rect(slide, Inches(x), Inches(1.55), Inches(2.9), Inches(4.9), RGBColor(0xF5, 0xF5, 0xF7))
        add_textbox(slide, Inches(x + 0.2), Inches(1.75), Inches(2.5), Inches(0.5),
                    icon, size=28)
        add_textbox(slide, Inches(x + 0.2), Inches(2.3), Inches(2.5), Inches(0.7),
                    title, size=15, bold=True, color=INK)
        add_textbox(slide, Inches(x + 0.2), Inches(3.1), Inches(2.5), Inches(2.5),
                    desc, size=12, color=MUTED)
        x += 3.1

    # ── 4. BRIEF HACKATHON ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Alignement brief hackathon", "100 % des exigences couvertes")
    rows = [
        ("Capture visuelle (obligatoire)", "Module /scan : caméra PWA, scan guidé 6 étapes, fallback import photo"),
        ("Gestion des pannes", "CRUD complet, priorités, commentaires, historique, images"),
        ("Pièces détachées", "Suivi commande → reçu → installé, lié aux pannes"),
        ("Workflow", "Kanban 7 colonnes avec notifications automatiques"),
        ("Multi-utilisateurs", "4 rôles : Admin, Manager, Technicien, Client"),
    ]
    y = 1.6
    for req, impl in rows:
        add_rect(slide, Inches(0.7), Inches(y), Inches(0.15), Inches(0.55), GREEN)
        add_textbox(slide, Inches(1.0), Inches(y), Inches(4.5), Inches(0.4),
                    req, size=14, bold=True, color=INK)
        add_textbox(slide, Inches(5.5), Inches(y), Inches(6.8), Inches(0.55),
                    impl, size=13, color=MUTED)
        y += 0.75

    # ── 5. PERSONAS ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "4 espaces métier", "Une interface adaptée à chaque acteur")
    personas = [
        ("Administrateur", "Centre de contrôle, utilisateurs, entreprises, audit, supervision globale", BRAND),
        ("Manager", "Dashboard KPI, Kanban, équipements, activité temps réel", RGBColor(0x0E, 0x7E, 0x96)),
        ("Technicien", "Mes interventions, calendrier, scan terrain, pièces", RGBColor(0x00, 0x71, 0xE3)),
        ("Client", "Portail dédié, déclaration panne, suivi demandes, onboarding", RGBColor(0x16, 0xA3, 0x4A)),
    ]
    x = 0.7
    for name, desc, col in personas:
        add_rect(slide, Inches(x), Inches(1.55), Inches(2.9), Inches(0.08), col)
        add_textbox(slide, Inches(x + 0.15), Inches(1.75), Inches(2.6), Inches(0.4),
                    name, size=16, bold=True, color=col)
        add_textbox(slide, Inches(x + 0.15), Inches(2.2), Inches(2.6), Inches(3.5),
                    desc, size=13, color=MUTED)
        x += 3.1

    # ── 6. CAPTURE VISUELLE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Module Capture Visuelle", "Cœur du brief — démo live recommandée")
    steps = ["1. Sélection équipement", "2. Scan guidé multi-angles", "3. Pré-diagnostic IA",
             "4. Déclaration panne", "5. Suivi Kanban"]
    add_bullets(slide, Inches(0.9), Inches(1.55), Inches(5.5), Inches(3), steps, size=16)
    add_bullets(slide, Inches(0.9), Inches(4.2), Inches(5.8), Inches(2.5), [
        "Réticule + ligne de balayage animée",
        "6 étapes : face, profils, plaque, zone défaut…",
        "Preuve visuelle obligatoire avant envoi",
        "Fallback : import photo si pas de caméra",
    ], size=14, color=MUTED)
    add_rect(slide, Inches(7), Inches(1.55), Inches(5.5), Inches(4.8), RGBColor(0x0F, 0x0F, 0x12))
    add_textbox(slide, Inches(7.3), Inches(1.85), Inches(5), Inches(0.5),
                "Parcours scan", size=18, bold=True, color=WHITE)
    flow = "Équipement → Caméra → Analyse → Panne créée → Kanban"
    add_textbox(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(1),
                flow, size=14, color=RGBColor(0xAE, 0xAE, 0xB2))
    add_textbox(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(1.2),
                "💡 Astuce démo : partir du Dashboard → Scanner → choisir « Compresseur à vis »",
                size=13, color=ACCENT)

    # ── 7. PRÉ-DIAGNOSTIC ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Innovation : Pré-diagnostic assisté", "L'IA qui accélère la première analyse")
    add_textbox(slide, Inches(0.9), Inches(1.55), Inches(11), Inches(0.6),
                "Heuristique intelligente combinant type d'équipement, health score et couverture du scan",
                size=15, color=MUTED)
    cards = [
        ("Anomalie suggérée", "Ex : « Fuite hydraulique », « Usure roulement convoyeur »"),
        ("Gravité auto", "Critical / High / Medium selon health score"),
        ("Pièces probables", "Joint SPI, capteur inductif, roulement…"),
        ("Score confiance", "Jusqu'à 96 % selon clichés + état machine"),
    ]
    x, y = 0.7, 2.3
    for i, (t, d) in enumerate(cards):
        if i == 2:
            x, y = 0.7, 4.5
        add_rect(slide, Inches(x), Inches(y), Inches(5.8), Inches(1.8), RGBColor(0xEE, 0xF6, 0xFF))
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.15), Inches(5.4), Inches(0.4),
                    t, size=15, bold=True, color=BRAND)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.55), Inches(5.4), Inches(1),
                    d, size=13, color=INK)
        x += 6.2

    # ── 8. WORKFLOW ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Workflow bout-en-bout", "De la panne à la pièce installée")
    workflow = "Soumis → Analyse → Inspection → Validation → Fabrication → Livraison → Clôturé"
    add_rect(slide, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.9), BRAND)
    add_textbox(slide, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.5),
                workflow, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(0.9), Inches(2.8), Inches(5.5), Inches(3.5), [
        "Kanban drag & drop (style Trello)",
        "Chaque déplacement = notification temps réel",
        "Fiche panne : commentaires, images, pièces liées",
        "Pièces : En attente → Commandé → Reçu → Installé",
    ], size=15)
    add_rect(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5), RGBColor(0xF5, 0xF5, 0xF7))
    add_textbox(slide, Inches(7.3), Inches(3.0), Inches(5), Inches(0.4),
                "Données de démo pré-chargées", size=14, bold=True, color=INK)
    add_bullets(slide, Inches(7.3), Inches(3.5), Inches(5), Inches(2.5), [
        "6 équipements industriels réalistes",
        "5 pannes à différents statuts",
        "3 pièces en cours de commande",
        "4 utilisateurs de test",
    ], size=13, color=MUTED)

    # ── 9. JUMEAU NUMÉRIQUE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Jumeau numérique des chantiers", "Vue 2D / 3D isométrique du parc")
    add_bullets(slide, Inches(0.9), Inches(1.6), Inches(6), Inches(4), [
        "Cartographie des bâtiments et zones (A, B, C, D)",
        "Équipements positionnés avec code couleur sévérité",
        "Pannes critiques visibles en un coup d'œil",
        "Clic → détail équipement ou lancement scan direct",
        "Projection isométrique 3D calculée en temps réel",
    ], size=15)
    add_rect(slide, Inches(7.2), Inches(1.6), Inches(5.2), Inches(4.5), RGBColor(0x0F, 0x0F, 0x12))
    add_textbox(slide, Inches(7.5), Inches(2.5), Inches(4.6), Inches(2),
                "Différenciant hackathon :\nvisualisation spatiale\n+ maintenance terrain",
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 10. DASHBOARD ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Pilotage & Health Score", "Décisions basées sur les données")
    kpis = [
        ("Health Score", "0–100 par équipement, calculé dynamiquement"),
        ("KPIs temps réel", "Pannes ouvertes, critiques, pièces en attente"),
        ("Graphiques", "Répartition par statut, tendances (Recharts)"),
        ("Admin", "Multi-entreprises, journal d'audit, utilisateurs en ligne"),
    ]
    x = 0.7
    for t, d in kpis:
        add_rect(slide, Inches(x), Inches(1.6), Inches(2.9), Inches(2.2), RGBColor(0xF5, 0xF5, 0xF7))
        add_textbox(slide, Inches(x + 0.2), Inches(1.8), Inches(2.5), Inches(0.5),
                    t, size=14, bold=True, color=BRAND)
        add_textbox(slide, Inches(x + 0.2), Inches(2.35), Inches(2.5), Inches(1.2),
                    d, size=12, color=MUTED)
        x += 3.1
    add_bullets(slide, Inches(0.9), Inches(4.2), Inches(11), Inches(2), [
        "PWA installable · Mode clair / sombre · Notifications activité",
        "Fonctionne offline (Service Worker + cache API)",
    ], size=14)

    # ── 11. ARCHITECTURE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Architecture technique", "Full-stack TypeScript, déployé en production")
    layers = [
        ("Frontend", "React 19 · TypeScript · Vite 8 · Tailwind · PWA"),
        ("Backend", "Node.js · Express · API REST · Auth JWT"),
        ("Data", "SQLite (sql.js) · Seed automatique · Migrations"),
        ("Deploy", "Netlify (frontend + serverless functions)"),
    ]
    y = 1.6
    for layer, tech in layers:
        add_rect(slide, Inches(0.7), Inches(y), Inches(2.2), Inches(0.65), BRAND)
        add_textbox(slide, Inches(0.85), Inches(y + 0.12), Inches(2), Inches(0.4),
                    layer, size=13, bold=True, color=WHITE)
        add_textbox(slide, Inches(3.1), Inches(y + 0.12), Inches(9.3), Inches(0.5),
                    tech, size=14, color=INK)
        y += 0.85
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(0.5),
                "Développé en 48h · Monorepo · 20+ pages · API complète · Démo live",
                size=14, bold=True, color=BRAND)

    # ── 12. DÉMO LIVE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0xF5, 0xF5, 0xF7))
    slide_header(slide, "Démo live — Script 3 minutes", "maintanxpro.netlify.app")
    script = [
        ("0:00", "Connexion admin@maintainx.com / demo1234"),
        ("0:30", "Dashboard → KPIs + équipements à risque"),
        ("1:00", "Scan → Compresseur → capture → pré-diagnostic"),
        ("1:45", "Kanban → glisser une panne vers « Inspection »"),
        ("2:15", "Pièces → commander une pièce détachée"),
        ("2:45", "Chantiers 3D → vue spatiale + clic équipement"),
    ]
    y = 1.55
    for time, action in script:
        add_rect(slide, Inches(0.7), Inches(y), Inches(1.1), Inches(0.55), BRAND)
        add_textbox(slide, Inches(0.75), Inches(y + 0.1), Inches(1), Inches(0.35),
                    time, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.0), Inches(y + 0.1), Inches(10), Inches(0.4),
                    action, size=15, color=INK)
        y += 0.72
    add_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.9), WHITE)
    add_textbox(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.6),
                "Comptes démo : Admin · Sophie (Manager) · Thomas (Technicien) · Marie (Client) — mot de passe : demo1234",
                size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # ── 13. IMPACT ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Valeur & impact", "Ce que MaintainX Pro apporte")
    impacts = [
        ("−40 %", "temps de qualification\ninitial d'une panne"),
        ("100 %", "des déclarations\navec preuve visuelle"),
        ("1 plateforme", "pour 4 profils métier\ndistincts"),
        ("0 install", "PWA accessible\nsur mobile terrain"),
    ]
    x = 0.7
    for val, lbl in impacts:
        add_rect(slide, Inches(x), Inches(1.7), Inches(2.9), Inches(3.5), RGBColor(0xEE, 0xF6, 0xFF))
        add_textbox(slide, Inches(x + 0.2), Inches(2.2), Inches(2.5), Inches(0.8),
                    val, size=36, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.2), Inches(3.2), Inches(2.5), Inches(1.5),
                    lbl, size=14, color=MUTED, align=PP_ALIGN.CENTER)
        x += 3.1
    add_textbox(slide, Inches(0.9), Inches(5.5), Inches(11), Inches(0.8),
                "Vision : intégrer un vrai modèle de vision par ordinateur + maintenance prédictive IoT",
                size=14, color=INK)

    # ── 14. ROADMAP ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    slide_header(slide, "Perspectives", "Au-delà du hackathon")
    roadmap = [
        ("Court terme", "Modèle CV réel (détection défauts), notifications push, export PDF rapports"),
        ("Moyen terme", "Connecteurs ERP/SAP, capteurs IoT → health score prédictif"),
        ("Long terme", "Marketplace pièces détachées, jumeau numérique BIM complet"),
    ]
    y = 1.7
    for phase, items in roadmap:
        add_textbox(slide, Inches(0.9), Inches(y), Inches(3), Inches(0.4),
                    phase, size=16, bold=True, color=BRAND)
        add_textbox(slide, Inches(3.5), Inches(y), Inches(8.5), Inches(0.5),
                    items, size=14, color=INK)
        y += 0.9

    # ── 15. CLOSING ──
    slide_closing(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
